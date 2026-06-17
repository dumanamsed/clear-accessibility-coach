import os
import json
from io import BytesIO
from flask import Flask, request, render_template, jsonify, Response
from werkzeug.utils import secure_filename

import config
from clear_analyzer.models import Finding
from clear_analyzer.citations import (
    STRAND_DEFINITIONS, PRESSBOOK_LINKS, FRAMEWORK_CITATION, STRAND_ORDER,
    MANUAL_CHECKS,
)
from clear_analyzer.pptx_rules import analyze_pptx
from clear_analyzer.docx_rules import analyze_docx
from clear_analyzer.pdf_rules import analyze_pdf
from clear_analyzer.html_rules import analyze_html, analyze_text
from clear_analyzer.readability import analyze_readability
from clear_analyzer.claude_review import run_claude_review

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = config.MAX_UPLOAD_SIZE_BYTES


@app.route("/")
def index():
    return render_template(
        "index.html",
        ctl_email=config.CTL_CONTACT_EMAIL,
        framework_citation=FRAMEWORK_CITATION,
    )


@app.route("/analyze", methods=["POST"])
def analyze():
    """Phase 1: deterministic rule checks. Returns instantly.

    The response includes an `ai_context` blob (document outline, alt text,
    body sample) that the client POSTs back to /ai-review for phase 2. Round-
    tripping through the client keeps the server stateless — no upload content
    is ever stored server-side, consistent with the privacy promise.
    """
    paste_content = request.form.get("paste_content", "").strip()
    paste_type = request.form.get("paste_type", "text")
    file = request.files.get("file")

    if file and file.filename:
        filename = secure_filename(file.filename)
        ext = os.path.splitext(filename)[1].lower()

        if ext not in config.ALLOWED_EXTENSIONS:
            return jsonify({"error": f"Unsupported file type: {ext}"}), 400

        file_bytes = file.stream.read()
        findings, outline, alt_texts, body_sample, file_type = _analyze_file(
            file_bytes, ext
        )
    elif paste_content:
        ext = ".html" if paste_type == "html" else ".txt"
        file_bytes = paste_content.encode("utf-8")
        findings, outline, alt_texts, body_sample, file_type = _analyze_file(
            file_bytes, ext
        )
        filename = "pasted_content" + ext
    else:
        return jsonify({"error": "No file or content provided."}), 400

    has_media = any(f.strand == "C" for f in findings)
    media_count = sum(1 for f in findings if f.strand == "C")

    report_data = _build_report_data(findings, has_media, media_count)
    report_data["claude_enabled"] = bool(config.ANTHROPIC_API_KEY)
    report_data["filename"] = filename if file and file.filename else "Pasted content"
    report_data["ai_context"] = {
        "file_type": file_type,
        "outline": outline,
        "alt_texts": alt_texts,
        "body_sample": body_sample,
    }

    return jsonify(report_data)


@app.route("/ai-review", methods=["POST"])
def ai_review():
    """Phase 2: Claude qualitative coaching pass.

    Called by the client after the rule report renders, with the ai_context
    from /analyze plus the rule findings (so Claude doesn't repeat them).
    """
    payload = request.get_json(silent=True) or {}
    ctx = payload.get("ai_context") or {}
    rule_dicts = payload.get("rule_findings") or []

    rule_findings = [
        Finding(
            strand=d.get("strand", ""),
            severity=d.get("severity", "tip"),
            location=d.get("location", ""),
            issue=d.get("issue", ""),
            evidence=d.get("evidence", ""),
        )
        for d in rule_dicts
        if isinstance(d, dict)
    ]

    claude_findings = run_claude_review(
        file_type=ctx.get("file_type", "document"),
        outline=ctx.get("outline") or [],
        alt_texts=ctx.get("alt_texts") or [],
        rule_findings=rule_findings,
        body_sample=ctx.get("body_sample", ""),
    )

    # None = the AI pass failed/was unavailable. An EMPTY list is a success:
    # the review ran and the document needed no extra coaching.
    return jsonify({
        "claude_available": claude_findings is not None,
        "findings": [f.to_dict() for f in (claude_findings or [])],
    })


@app.route("/export-pdf", methods=["POST"])
def export_pdf():
    report_data = request.get_json()
    if not report_data:
        return jsonify({"error": "No report data provided."}), 400

    html_content = render_template(
        "report.html",
        report=report_data,
        strand_definitions=STRAND_DEFINITIONS,
        strand_order=STRAND_ORDER,
        pressbook_links=PRESSBOOK_LINKS,
        manual_checks=MANUAL_CHECKS,
        framework_citation=FRAMEWORK_CITATION,
        ctl_email=config.CTL_CONTACT_EMAIL,
        is_pdf=True,
    )

    # Prefer a true server-side PDF (WeasyPrint works on Linux hosts like
    # Render). The report must model what it preaches, so we request the
    # PDF/UA-1 variant — a TAGGED, screen-reader-navigable PDF with the
    # document language and title set. If the installed WeasyPrint doesn't
    # support the variant, fall back to a plain PDF; if its native libraries
    # are missing entirely (e.g. macOS without Homebrew), fall back to the
    # print-styled HTML — the client detects the content type and opens the
    # browser print dialog instead.
    try:
        from weasyprint import HTML as WeasyHTML

        doc = WeasyHTML(string=html_content)
        try:
            pdf_bytes = doc.write_pdf(pdf_variant="pdf/ua-1")
        except Exception:
            pdf_bytes = doc.write_pdf()
        return Response(
            pdf_bytes,
            mimetype="application/pdf",
            headers={
                "Content-Disposition": "attachment; filename=clear-accessibility-report.pdf"
            },
        )
    except Exception as exc:
        import sys
        print(
            f"[export_pdf] WeasyPrint unavailable, falling back to print view: "
            f"{type(exc).__name__}",
            file=sys.stderr,
            flush=True,
        )
        return html_content


def _analyze_file(file_bytes, ext):
    findings = []
    outline = []
    alt_texts = []
    body_sample = ""
    file_type = ext.lstrip(".")

    if ext == ".pptx":
        findings = analyze_pptx(file_bytes)
        outline, alt_texts, body_sample = _extract_pptx_metadata(file_bytes)
    elif ext == ".docx":
        findings = analyze_docx(file_bytes)
        outline, alt_texts, body_sample = _extract_docx_metadata(file_bytes)
    elif ext == ".pdf":
        findings = analyze_pdf(file_bytes)
        outline, alt_texts, body_sample = _extract_pdf_metadata(file_bytes)
    elif ext in (".html", ".htm"):
        content = file_bytes.decode("utf-8", errors="replace")
        findings = analyze_html(content)
        outline, alt_texts, body_sample = _extract_html_metadata(content)
        file_type = "html"
    elif ext == ".txt":
        content = file_bytes.decode("utf-8", errors="replace")
        findings = analyze_text(content)
        body_sample = content[:2000]
        file_type = "text"

    # CLEAR "Easy to Read" readability scoring (Flesch / sentence length) runs on
    # the extracted prose, independent of file type.
    findings = findings + analyze_readability(body_sample)

    return findings, outline, alt_texts, body_sample, file_type


def _extract_pptx_metadata(file_bytes):
    from pptx import Presentation
    prs = Presentation(BytesIO(file_bytes))
    outline = []
    alt_texts = []
    body_parts = []

    for i, slide in enumerate(prs.slides, 1):
        title = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text_frame.text.strip()
        outline.append(f"Slide {i}: {title or '(no title)'}")

        for shape in slide.shapes:
            if hasattr(shape, '_element'):
                for elem in shape._element.iter():
                    if elem.tag.endswith("}cNvPr"):
                        descr = elem.get("descr", "")
                        if descr and descr.strip():
                            alt_texts.append(f"Slide {i}: {descr.strip()}")

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        body_parts.append(para.text.strip())

    return outline, alt_texts, "\n".join(body_parts)[:2000]


def _extract_docx_metadata(file_bytes):
    from docx import Document
    from docx.oxml.ns import qn
    doc = Document(BytesIO(file_bytes))
    outline = []
    alt_texts = []
    body_parts = []

    for para in doc.paragraphs:
        style_name = para.style.name if para.style else ""
        if style_name.startswith("Heading"):
            outline.append(f"{style_name}: {para.text.strip()}")

        if para.text.strip():
            body_parts.append(para.text.strip())

        for run in para.runs:
            for docPr in run._element.findall(f".//{qn('wp:docPr')}"):
                descr = docPr.get("descr", "")
                if descr and descr.strip():
                    alt_texts.append(descr.strip())

    return outline, alt_texts, "\n".join(body_parts)[:2000]


def _extract_pdf_metadata(file_bytes):
    import fitz
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    outline = []
    body_parts = []

    toc = doc.get_toc()
    for level, title, page in toc:
        outline.append(f"{'  ' * (level - 1)}{title} (p.{page})")

    for page in doc:
        text = page.get_text("text").strip()
        if text:
            body_parts.append(text)

    doc.close()
    return outline, [], "\n".join(body_parts)[:2000]


def _extract_html_metadata(content):
    import re
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(content, "html.parser")
    outline = []
    alt_texts = []

    for h in soup.find_all(re.compile(r"^h[1-6]$")):
        outline.append(f"<{h.name}>: {h.get_text(strip=True)}")

    for img in soup.find_all("img"):
        alt = img.get("alt", "")
        if alt and alt.strip():
            alt_texts.append(alt.strip())

    body_text = soup.get_text(separator="\n", strip=True)
    return outline, alt_texts, body_text[:2000]


def _build_report_data(findings, has_media, media_count):
    strands = {}
    for strand_key in STRAND_ORDER:
        strand_findings = [f for f in findings if f.strand == strand_key]
        critical = sum(1 for f in strand_findings if f.severity == "critical")
        warnings = sum(1 for f in strand_findings if f.severity == "warning")
        tips = sum(1 for f in strand_findings if f.severity == "tip")
        strands[strand_key] = {
            "findings": [f.to_dict() for f in strand_findings],
            "critical": critical,
            "warning": warnings,
            "tip": tips,
            "total": len(strand_findings),
        }

    total_findings = sum(s["total"] for s in strands.values())
    strands_with_findings = sum(1 for s in strands.values() if s["total"] > 0)

    return {
        "strands": strands,
        "total_findings": total_findings,
        "strands_with_findings": strands_with_findings,
        "has_media": has_media,
        "media_count": media_count,
    }


if __name__ == "__main__":
    app.run(host="127.0.0.1", port=config.FLASK_PORT, debug=True)
