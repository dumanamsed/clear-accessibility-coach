import re
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from .models import Finding
from .contrast import office_rgb, contrast_ratio, required_ratio, is_near

# Common image file extensions
_IMAGE_EXT_RE = re.compile(
    r"^.+\.(png|jpg|jpeg|gif|bmp|tiff|tif|svg|webp|ico|emf|wmf)$",
    re.IGNORECASE,
)

# Patterns that indicate auto-generated or meaningless alt text
_WEAK_ALT_PATTERNS = [
    re.compile(r"^image\s*\d*$", re.IGNORECASE),          # "image", "Image1"
    re.compile(r"^picture\s*\d*$", re.IGNORECASE),         # "picture", "Picture 2"
    re.compile(r"^photo\s*\d*$", re.IGNORECASE),           # "photo"
    re.compile(r"^img[_\s]?\d*$", re.IGNORECASE),          # "img_1", "img 2"
    re.compile(r"^screenshot", re.IGNORECASE),             # "screenshot..."
    re.compile(r"^graphic\s*\d*$", re.IGNORECASE),         # "graphic"
    re.compile(r"^chart\s*\d*$", re.IGNORECASE),           # "chart 1"
    re.compile(r"^figure\s*\d*$", re.IGNORECASE),          # "figure"
    re.compile(r"^slide\s*\d+", re.IGNORECASE),            # "Slide 1"
    re.compile(r"^content placeholder", re.IGNORECASE),    # PowerPoint default
]


def analyze_pptx(file_bytes: bytes) -> list[Finding]:
    findings = []
    prs = Presentation(BytesIO(file_bytes))

    for slide_num, slide in enumerate(prs.slides, 1):
        loc = f"Slide {slide_num}"
        _check_images(slide, loc, findings)
        _check_title(slide, loc, findings)
        _check_media(slide, loc, findings)
        _check_text(slide, loc, findings)
        _check_reading_order(slide, loc, findings)
        _check_justified(slide, loc, findings)
        _check_text_contrast(slide, loc, findings)

    return findings


def _solid_rgb(fill):
    """Return an explicit solid-fill RGB for a fill, or None if it's themed,
    gradient, picture, or otherwise not a concrete RGB we can evaluate."""
    try:
        from pptx.enum.dml import MSO_FILL
        if fill.type != MSO_FILL.SOLID:
            return None
        return office_rgb(fill.fore_color.rgb)
    except Exception:
        return None


def _slide_bg_rgb(slide):
    """Best-effort: the slide's own solid background fill, if explicitly set."""
    try:
        bg = slide.background
        return _solid_rgb(bg.fill)
    except Exception:
        return None


def _check_text_contrast(slide, loc, findings):
    """WCAG 1.4.3 for slides. Contrast on slides is genuinely hard because
    backgrounds are often theme- or layout-driven; to stay false-positive-free
    we ONLY evaluate a run when BOTH its font color AND the background behind it
    (the shape's own solid fill, else the slide's solid background) are explicit
    RGB values. Themed colors are skipped rather than guessed."""
    slide_bg = _slide_bg_rgb(slide)
    flagged = 0
    for shape in slide.shapes:
        if flagged >= 3 or not getattr(shape, "has_text_frame", False):
            continue
        shape_bg = _solid_rgb(shape.fill) if hasattr(shape, "fill") else None
        bg = shape_bg or slide_bg
        if bg is None:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                try:
                    fg = office_rgb(run.font.color.rgb) if run.font.color and run.font.color.type is not None else None
                except Exception:
                    fg = None
                if fg is None:
                    continue
                size = run.font.size.pt if run.font.size else None
                font_px = size * 1.333 if size else None
                ratio = contrast_ratio(fg, bg)
                needed = required_ratio(font_px, bool(run.font.bold))
                if ratio + 0.05 < needed:
                    findings.append(Finding(
                        strand="E",
                        severity="warning",
                        location=loc,
                        issue=f"Text color #%02X%02X%02X has only {ratio:.1f}:1 contrast against its "
                              f"background #%02X%02X%02X, below the WCAG 2.2 AA minimum of {needed:g}:1 "
                              f"(Success Criterion 1.4.3)." % (fg[0], fg[1], fg[2], bg[0], bg[1], bg[2]),
                        evidence=run.text.strip()[:80],
                    ))
                    flagged += 1
                    break
            if flagged >= 3:
                break


def _check_justified(slide, loc, findings):
    """CLEAR Easy to Read: prefer left-aligned over justified text."""
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for para in shape.text_frame.paragraphs:
            if para.alignment == PP_ALIGN.JUSTIFY and len(para.text.split()) > 10:
                findings.append(Finding(
                    strand="E",
                    severity="tip",
                    location=loc,
                    issue="Text is justified, which creates uneven word spacing that is harder "
                          "to read. CLEAR recommends left-aligned text.",
                    evidence=para.text.strip()[:80],
                ))
                return  # one per slide is enough


def _check_reading_order(slide, loc, findings):
    """Heuristic: floating text boxes (added outside the layout's placeholders)
    are read by screen readers in insertion order, which often differs from the
    visual order. Two or more on a slide is a common reading-order problem."""
    floating = 0
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if not shape.text_frame.text.strip():
            continue
        try:
            is_placeholder = shape.placeholder_format is not None
        except ValueError:
            is_placeholder = False
        if not is_placeholder:
            floating += 1

    if floating >= 2:
        findings.append(Finding(
            strand="L",
            severity="warning",
            location=loc,
            issue=f"Slide has {floating} floating text boxes. Screen readers read them in the order they were added, which may not match the visual layout.",
            evidence="Use the layout's placeholders where possible, and check Home > Arrange > Selection Pane for reading order.",
        ))


def _get_alt_text(shape):
    """Extract alt text from any shape by walking the XML tree.

    python-pptx doesn't expose alt text directly on all shape types.
    Alt text lives in the 'descr' attribute of the cNvPr element inside
    the shape's non-visual properties (nvSpPr, nvPicPr, nvGrpSpPr, etc.).
    We also check the newer a16:creationId parent's descr.
    """
    el = shape._element

    # Method 1: Find any cNvPr element anywhere in this shape's XML
    for cNvPr in el.iter():
        if cNvPr.tag.endswith("}cNvPr") or cNvPr.tag == "cNvPr":
            descr = cNvPr.get("descr", "")
            if descr and descr.strip():
                return descr.strip()

    # Method 2: Check for decorative flag (Office 365+)
    # <adec:decorative val="1"/> means intentionally decorative
    for dec in el.iter():
        if "decorative" in dec.tag.lower():
            val = dec.get("val", "")
            if val == "1":
                return "__DECORATIVE__"

    return ""


def _check_images(slide, loc, findings):
    """Check ALL shapes that contain images, not just PICTURE type."""
    seen_shapes = set()

    for shape in slide.shapes:
        _check_shape_for_image(shape, loc, findings, seen_shapes)


def _check_shape_for_image(shape, loc, findings, seen_shapes):
    """Recursively check a shape (and grouped shapes) for images."""
    shape_id = id(shape)
    if shape_id in seen_shapes:
        return
    seen_shapes.add(shape_id)

    # Check grouped shapes recursively
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        if hasattr(shape, "shapes"):
            for child_shape in shape.shapes:
                _check_shape_for_image(child_shape, loc, findings, seen_shapes)
        return

    # Determine if this shape contains an image
    is_image = False

    # Direct picture type
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        is_image = True

    # Placeholder with an image fill
    if not is_image:
      try:
        is_placeholder = shape.placeholder_format is not None
      except ValueError:
        is_placeholder = False
      if is_placeholder:
        el = shape._element
        blipFills = list(el.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}blipFill"))
        if blipFills:
            is_image = True

    # Any shape with a blipFill (image fill) that isn't a placeholder
    if not is_image:
        el = shape._element
        blipFills = list(el.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}blipFill"))
        if blipFills:
            is_image = True

    # Shape with an image in its spPr (shape properties) via blipFill
    if not is_image:
        el = shape._element
        # Check for pic element
        pics = list(el.iter("{http://schemas.openxmlformats.org/presentationml/2006/main}pic"))
        if not pics:
            pics = list(el.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}pic"))
        if pics:
            is_image = True

    if not is_image:
        return

    alt_text = _get_alt_text(shape)

    if alt_text == "__DECORATIVE__":
        # Marked as decorative — this is correct, no finding needed
        return

    name = shape.name or "Image"

    if not alt_text:
        findings.append(Finding(
            strand="A",
            severity="critical",
            location=loc,
            issue="Image is missing alt text.",
            evidence=f"Shape: {name}",
        ))
    elif _IMAGE_EXT_RE.match(alt_text):
        # Alt text is just a filename like "IMG_4032.jpg"
        findings.append(Finding(
            strand="A",
            severity="critical",
            location=loc,
            issue="Image alt text is a filename, not a meaningful description.",
            evidence=f"Shape: {name}, alt text: \"{alt_text}\"",
        ))
    elif any(p.match(alt_text) for p in _WEAK_ALT_PATTERNS):
        # Alt text is generic/auto-generated
        findings.append(Finding(
            strand="A",
            severity="warning",
            location=loc,
            issue="Image alt text appears to be generic or auto-generated. Consider a description of what the image conveys.",
            evidence=f"Shape: {name}, alt text: \"{alt_text}\"",
        ))


def _check_title(slide, loc, findings):
    has_title = False
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        title_text = slide.shapes.title.text_frame.text.strip()
        if title_text:
            has_title = True

    if not has_title:
        for shape in slide.shapes:
            try:
                pf = shape.placeholder_format
            except ValueError:
                continue
            if shape.has_text_frame and pf is not None:
                if pf.idx == 0 and shape.text_frame.text.strip():
                    has_title = True
                    break

    if not has_title:
        findings.append(Finding(
            strand="L",
            severity="warning",
            location=loc,
            issue="Slide does not have a filled title.",
            evidence="Title placeholder is empty or missing.",
        ))


def _check_media(slide, loc, findings):
    detected = set()
    for shape in slide.shapes:
        if shape.shape_type in (
            MSO_SHAPE_TYPE.MEDIA,
            MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
        ):
            detected.add(shape.name or "media object")
            findings.append(Finding(
                strand="C",
                severity="warning",
                location=loc,
                issue="Embedded media detected. Confirm that captions or a transcript are available.",
                evidence=f"Shape: {shape.name or 'media object'}",
            ))
            continue

        # Check XML for video/audio references not caught by shape_type
        if hasattr(shape, '_element'):
            el = shape._element
            el_xml = el.xml if hasattr(el, 'xml') else ""
            shape_name = shape.name or "media element"
            if shape_name not in detected:
                if "video" in el_xml.lower() or "audio" in el_xml.lower():
                    # Avoid false positives from shape names containing "video"
                    has_media_ref = False
                    for tag in ("videoFile", "audioFile", "MediaBookmark"):
                        if tag.lower() in el_xml.lower():
                            has_media_ref = True
                            break
                    # Also check for OLE or embedded media relationships
                    if not has_media_ref:
                        for tag in ("<a:videoFile", "<a:audioFile", "r:link", "r:embed"):
                            if tag in el_xml:
                                has_media_ref = True
                                break
                    if has_media_ref:
                        detected.add(shape_name)
                        findings.append(Finding(
                            strand="C",
                            severity="warning",
                            location=loc,
                            issue="Embedded media detected. Confirm that captions or a transcript are available.",
                            evidence=f"Shape: {shape_name}",
                        ))


def _check_text(slide, loc, findings):
    small_font_flagged = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for para in shape.text_frame.paragraphs:
            if not small_font_flagged:
                for run in para.runs:
                    font = run.font
                    if font.size and font.size < Pt(18):
                        findings.append(Finding(
                            strand="E",
                            severity="tip",
                            location=loc,
                            issue=f"Text uses a small font size ({font.size.pt:.0f}pt). Consider 18pt or larger for readability.",
                            evidence=run.text[:80] if run.text else "(empty run)",
                        ))
                        small_font_flagged = True
                        break

            for run in para.runs:
                if run.hyperlink and run.hyperlink.address:
                    link_text = run.text.strip().lower()
                    if link_text in ("click here", "here", "link") or (
                        link_text.startswith("http://") or link_text.startswith("https://")
                    ):
                        findings.append(Finding(
                            strand="E",
                            severity="tip",
                            location=loc,
                            issue="Hyperlink uses generic or bare-URL text.",
                            evidence=f'Link text: "{run.text.strip()}"',
                        ))

        try:
            pf = shape.placeholder_format
        except ValueError:
            pf = None
        if shape.has_text_frame and pf is not None:
            idx = pf.idx
            if idx not in (0, 1):
                continue
            bullet_count = sum(
                1 for p in shape.text_frame.paragraphs if p.text.strip()
            )
            if bullet_count > 7:
                findings.append(Finding(
                    strand="E",
                    severity="tip",
                    location=loc,
                    issue=f"Text frame has {bullet_count} bullet points. Consider limiting to 7 or fewer per slide.",
                    evidence=f"Shape: {shape.name or 'text frame'}",
                ))
